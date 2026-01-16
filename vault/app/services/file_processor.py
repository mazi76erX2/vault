"""
File processing utilities for knowledge base ingestion
Supports: PDF, DOCX, PPTX, TXT, MD, CSV, XLSX, HTML, images (OCR), audio
"""
import logging
from pathlib import Path
from typing import List

logger = logging.getLogger(__name__)


def process_file(file_path: Path, file_ext: str) -> str:
    """
    Extract text from various file formats.
    
    Args:
        file_path: Path to the file
        file_ext: File extension (e.g., '.pdf', '.docx')
        
    Returns:
        Extracted text content
    """
    file_ext = file_ext.lower()
    
    try:
        if file_ext == ".pdf":
            return extract_from_pdf(file_path)
        elif file_ext in [".docx", ".doc"]:
            return extract_from_docx(file_path)
        elif file_ext == ".pptx":
            return extract_from_pptx(file_path)
        elif file_ext in [".txt", ".md", ".markdown"]:
            return extract_from_text(file_path)
        elif file_ext == ".csv":
            return extract_from_csv(file_path)
        elif file_ext in [".xlsx", ".xls"]:
            return extract_from_excel(file_path)
        elif file_ext in [".html", ".htm"]:
            return extract_from_html(file_path)
        elif file_ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
            return extract_from_image_ocr(file_path)
        elif file_ext in [".mp3", ".wav", ".m4a", ".ogg", ".webm"]:
            return extract_from_audio(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
            
    except Exception as e:
        logger.error(f"Error processing {file_ext} file: {e}")
        raise


def extract_from_pdf(file_path: Path) -> str:
    """Extract text from PDF using PyPDF2."""
    try:
        from PyPDF2 import PdfReader
        
        reader = PdfReader(str(file_path))
        text = []
        
        for page_num, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    text.append(f"--- Page {page_num + 1} ---\n{page_text}")
            except Exception as e:
                logger.warning(f"Error extracting page {page_num + 1}: {e}")
                continue
        
        return "\n\n".join(text)
        
    except ImportError:
        raise ImportError("PyPDF2 required: pip install PyPDF2")


def extract_from_docx(file_path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        
        doc = Document(str(file_path))
        text = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text.append(row_text)
        
        return "\n\n".join(text)
        
    except ImportError:
        raise ImportError("python-docx required: pip install python-docx")


def extract_from_pptx(file_path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        
        prs = Presentation(str(file_path))
        text = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = [f"--- Slide {slide_num + 1} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:
                text.append("\n".join(slide_text))
        
        return "\n\n".join(text)
        
    except ImportError:
        raise ImportError("python-pptx required: pip install python-pptx")


def extract_from_text(file_path: Path) -> str:
    """Extract text from TXT/MD files."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_from_csv(file_path: Path) -> str:
    """Extract text from CSV files."""
    import csv
    
    text = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = " | ".join(str(cell) for cell in row if cell)
            if row_text.strip():
                text.append(row_text)
    
    return "\n".join(text)


def extract_from_excel(file_path: Path) -> str:
    """Extract text from Excel files."""
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"=== Sheet: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text.append(row_text)
        
        return "\n\n".join(text)
        
    except ImportError:
        raise ImportError("openpyxl required: pip install openpyxl")


def extract_from_html(file_path: Path) -> str:
    """Extract text from HTML files."""
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)
        
        return text
        
    except ImportError:
        raise ImportError("beautifulsoup4 required: pip install beautifulsoup4")


def extract_from_image_ocr(file_path: Path) -> str:
    """Extract text from images using Tesseract OCR."""
    try:
        from PIL import Image
        import pytesseract
        
        image = Image.open(str(file_path))
        text = pytesseract.image_to_string(image)
        
        if not text.strip():
            raise ValueError("No text detected in image")
        
        return text
        
    except ImportError:
        raise ImportError(
            "Pillow and pytesseract required: pip install Pillow pytesseract\n"
            "Also install Tesseract: brew install tesseract (macOS) or apt-get install tesseract-ocr (Ubuntu)"
        )


def extract_from_audio(file_path: Path) -> str:
    """
    Transcribe audio using Faster Whisper (no PyTorch dependency).
    """
    try:
        from faster_whisper import WhisperModel
        
        # Use base model on CPU with int8 for better compatibility
        model = WhisperModel("base", device="cpu", compute_type="int8")
        
        segments, info = model.transcribe(str(file_path))
        
        # Combine all segments
        text = " ".join([segment.text for segment in segments])
        
        if not text.strip():
            raise ValueError("No speech detected in audio")
        
        return text
        
    except ImportError:
        raise ImportError(
            "faster-whisper not installed.\n"
            "Install with: pip install faster-whisper\n"
            "Also requires ffmpeg: brew install ffmpeg"
        )
    except Exception as e:
        logger.error(f"Audio transcription failed: {e}")
        raise ValueError(f"Failed to transcribe audio: {str(e)}")



def scrape_website(url: str) -> str:
    """
    Scrape text content from a website.
    
    Args:
        url: Website URL to scrape
        
    Returns:
        Extracted text content
    """
    try:
        import requests
        from bs4 import BeautifulSoup
        
        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"
        }
        
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, "html.parser")
        
        # Remove unwanted elements
        for element in soup(["script", "style", "nav", "footer", "header", "aside"]):
            element.decompose()
        
        # Extract text
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)
        
        # Add URL metadata
        title = soup.find("title")
        title_text = title.get_text() if title else "No Title"
        
        return f"# {title_text}\nSource: {url}\n\n{text}"
        
    except Exception as e:
        logger.error(f"Error scraping {url}: {e}")
        raise


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 100) -> List[str]:
    """
    Split text into overlapping chunks for better retrieval.
    
    Args:
        text: Text to chunk
        chunk_size: Target size of each chunk (characters)
        overlap: Overlap between chunks (characters)
        
    Returns:
        List of text chunks
    """
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        
        if end < len(text):
            for sep in [". ", ".\n", "! ", "?\n", "? "]:
                last_sep = text[start:end].rfind(sep)
                if last_sep != -1:
                    end = start + last_sep + len(sep)
                    break
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        start = end - overlap
    
    return chunks